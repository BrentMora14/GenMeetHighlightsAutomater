"""
UP CURSOR GenMeet Highlights Automater
  Tab 1 — Canva PPTX  →  Google Doc (formatted outline)
  Tab 2 — Google Doc  →  Canva Bulk Create CSV
  Tab 3 — F2F + Zoom CSVs  →  Consolidated attendance
"""

import csv
import io
import json
import re
import time
from typing import Optional

import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import attendance as att

# ══════════════════════════════════════════════════════════════════════════════
# PAGE CONFIG & CUSTOM CSS
# ══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="UP CURSOR GenMeet Highlights Automater",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="collapsed",
)

st.markdown("""
<style>
  @import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@300;400;500;600&family=DM+Mono:wght@400;500&display=swap');

  html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }

  /* ── Top header bar ── */
  .top-header {
    background: linear-gradient(135deg, #0f0f23 0%, #1a1a3e 60%, #0d1b2a 100%);
    padding: 2rem 2.5rem 1.5rem;
    border-radius: 16px;
    margin-bottom: 1.5rem;
    border: 1px solid rgba(255,255,255,0.08);
  }
  .top-header h1 {
    color: #ffffff;
    font-size: 1.9rem;
    font-weight: 600;
    margin: 0 0 .35rem 0;
    letter-spacing: -.02em;
  }
  .top-header p {
    color: rgba(255,255,255,0.55);
    font-size: .9rem;
    margin: 0;
  }
  .badge {
    display: inline-block;
    background: rgba(99,102,241,0.25);
    color: #a5b4fc;
    font-size: .7rem;
    font-weight: 600;
    letter-spacing: .08em;
    text-transform: uppercase;
    padding: .2rem .6rem;
    border-radius: 100px;
    border: 1px solid rgba(99,102,241,0.4);
    margin-right: .5rem;
    margin-bottom: .6rem;
  }

  /* ── Cards ── */
  .card {
    background: #fafafa;
    border: 1px solid #e8e8ee;
    border-radius: 12px;
    padding: 1.4rem 1.6rem;
    margin-bottom: 1rem;
  }
  .card-title {
    font-size: .8rem;
    font-weight: 600;
    letter-spacing: .07em;
    text-transform: uppercase;
    color: #6366f1;
    margin-bottom: .5rem;
  }

  /* ── Outline preview ── */
  .outline-box {
    background: #0f0f1a;
    color: #e2e8f0;
    font-family: 'DM Mono', monospace;
    font-size: .8rem;
    line-height: 1.8;
    padding: 1.2rem 1.4rem;
    border-radius: 10px;
    white-space: pre;
    overflow-x: auto;
    max-height: 380px;
    overflow-y: auto;
    border: 1px solid rgba(99,102,241,0.3);
  }
  .outline-box .lvl0 { color: #93c5fd; font-weight: 600; }
  .outline-box .lvl1 { color: #c4b5fd; }
  .outline-box .lvl2 { color: #6ee7b7; }
  .outline-box .lvl3 { color: #fde68a; }
  .outline-box .lvl4 { color: #fca5a5; }

  /* ── Step pills ── */
  .step-pill {
    display: inline-flex;
    align-items: center;
    gap: .4rem;
    background: #f0f0ff;
    color: #4338ca;
    font-size: .78rem;
    font-weight: 600;
    padding: .3rem .8rem;
    border-radius: 100px;
    border: 1px solid #c7d2fe;
    margin-bottom: 1rem;
  }

  /* ── Info/success boxes ── */
  .info-box {
    background: #eff6ff;
    border: 1px solid #bfdbfe;
    border-radius: 10px;
    padding: .9rem 1.1rem;
    color: #1e40af;
    font-size: .85rem;
    margin: .8rem 0;
  }
  .success-box {
    background: #f0fdf4;
    border: 1px solid #86efac;
    border-radius: 10px;
    padding: .9rem 1.1rem;
    color: #15803d;
    font-size: .85rem;
    margin: .8rem 0;
  }
  .warn-box {
    background: #fffbeb;
    border: 1px solid #fcd34d;
    border-radius: 10px;
    padding: .9rem 1.1rem;
    color: #92400e;
    font-size: .85rem;
    margin: .8rem 0;
  }

  /* ── Canva field table ── */
  .field-table { width: 100%; border-collapse: collapse; font-size: .85rem; }
  .field-table th {
    background: #f5f3ff;
    color: #4338ca;
    font-weight: 600;
    padding: .5rem .8rem;
    text-align: left;
    border-bottom: 2px solid #e0d9ff;
  }
  .field-table td {
    padding: .45rem .8rem;
    border-bottom: 1px solid #f0f0f0;
    color: #374151;
  }
  .field-table tr:last-child td { border-bottom: none; }
  .code-chip {
    background: #1e1e2e;
    color: #a5b4fc;
    font-family: 'DM Mono', monospace;
    font-size: .78rem;
    padding: .15rem .5rem;
    border-radius: 5px;
  }

  /* ── Streamlit overrides ── */
  div[data-testid="stFileUploader"] { border-radius: 10px; }
  .stButton > button {
    font-family: 'DM Sans', sans-serif;
    font-weight: 500;
    border-radius: 8px;
  }
  div[data-baseweb="tab"] button { font-family: 'DM Sans', sans-serif; font-weight: 500; }
  div[data-testid="stDataFrame"] { border-radius: 10px; overflow: hidden; }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# OUTLINE LOGIC (shared between tabs)
# ══════════════════════════════════════════════════════════════════════════════

def to_roman(n: int) -> str:
    val  = [1000,900,500,400,100,90,50,40,10,9,5,4,1]
    syms = ["M","CM","D","CD","C","XC","L","XL","X","IX","V","IV","I"]
    result = ""
    for v, s in zip(val, syms):
        while n >= v:
            result += s; n -= v
    return result

def to_alpha(n: int, upper: bool = True) -> str:
    return chr((64 if upper else 96) + n)

OUTLINE_PREFIX = {
    0: lambda n: f"{to_roman(n)}.",
    1: lambda n: f"{to_alpha(n)}.",
    2: lambda n: f"{n}.",
    3: lambda n: f"{to_alpha(n, False)})",
    4: lambda n: f"({n})",
}

INDENT = "    "   # 4 spaces per level


def _is_title_ph(shape) -> bool:
    if not shape.is_placeholder:
        return False
    try:
        ph = shape.placeholder_format.type
        return ph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False

def _avg_size(paras: list) -> float:
    sizes = [s for p in paras for s in p.get("_sizes", [])]
    return sum(sizes) / len(sizes) if sizes else 0


def extract_slides(pptx_bytes: bytes) -> list[dict]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides_out = []
    for slide in prs.slides:
        shapes_info = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            paras = []
            for p in shape.text_frame.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                sizes = [r.font.size.pt for r in p.runs if r.font and r.font.size]
                paras.append({"text": text, "level": p.level or 0, "_sizes": sizes})
            if not paras:
                continue
            shapes_info.append({
                "is_title": _is_title_ph(shape),
                "paras": paras,
                "top": shape.top or 0,
                "left": shape.left or 0,
                "avg_size": _avg_size(paras),
            })
        shapes_info.sort(key=lambda s: (s["top"], s["left"]))

        title_shape = next((s for s in shapes_info if s["is_title"]), None)
        if title_shape is None and shapes_info:
            title_shape = max(shapes_info, key=lambda s: s["avg_size"])

        title_text = ""
        body_shapes = []
        for s in shapes_info:
            if s is title_shape:
                title_text = " ".join(p["text"] for p in s["paras"])
            else:
                body_shapes.append(s)

        body = [{"text": p["text"], "level": p["level"]} for s in body_shapes for p in s["paras"]]
        slides_out.append({"title": title_text, "body": body})
    return slides_out


def build_numbered_outline(slides: list[dict]) -> list[dict]:
    """Returns list of {level, text, formatted} dicts.
    formatted uses tab characters for indentation; no outline prefixes.
    """
    raw = []
    for slide in slides:
        if slide["title"]:
            raw.append({"level": 0, "text": slide["title"]})
        for p in slide["body"]:
            raw.append({"level": p["level"] + 1, "text": p["text"]})

    return [
        {
            "level":     item["level"],
            "text":      item["text"],
            "formatted": "\t" * item["level"] + item["text"],
        }
        for item in raw
    ]


def outline_to_plain_text(outline: list[dict]) -> str:
    return "\n".join(item["formatted"] for item in outline)


def outline_to_html_preview(outline: list[dict]) -> str:
    css_cls = ["lvl0", "lvl1", "lvl2", "lvl3", "lvl4"]
    lines = []
    for item in outline:
        cls = css_cls[min(item["level"], 4)]
        # Tabs collapse in HTML — swap each for 4 spaces in the monospace preview
        display = item["formatted"].replace("\t", "    ")
        escaped = display.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        lines.append(f'<span class="{cls}">{escaped}</span>')
    return '<div class="outline-box">' + "\n".join(lines) + "</div>"


_BULLET_CHARS = ["•", "◦", "▪", "‣", "·"]


def _add_bullet_numbering(doc, max_level: int, num_id: int = 100, abstract_id: int = 100) -> int:
    """
    Inject a single multilevel bullet-list definition (ilvl 0..max_level)
    into the document's existing (initially empty) numbering part, and
    return the numId to use when attaching bullets to paragraphs.

    Word/Google Docs both represent outline nesting this way — a single
    numId whose ilvl on each paragraph *is* the outline level — which is
    exactly the convention read back by _docx_paragraph_level().
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    numbering_el = doc.part.numbering_part.element

    abstract_num = OxmlElement("w:abstractNum")
    abstract_num.set(qn("w:abstractNumId"), str(abstract_id))

    for i in range(max_level + 1):
        lvl = OxmlElement("w:lvl")
        lvl.set(qn("w:ilvl"), str(i))

        start = OxmlElement("w:start")
        start.set(qn("w:val"), "1")
        num_fmt = OxmlElement("w:numFmt")
        num_fmt.set(qn("w:val"), "bullet")
        lvl_text = OxmlElement("w:lvlText")
        lvl_text.set(qn("w:val"), _BULLET_CHARS[i % len(_BULLET_CHARS)])
        lvl_jc = OxmlElement("w:lvlJc")
        lvl_jc.set(qn("w:val"), "left")

        p_pr = OxmlElement("w:pPr")
        ind = OxmlElement("w:ind")
        ind.set(qn("w:left"), str(360 * (i + 1)))
        ind.set(qn("w:hanging"), "360")
        p_pr.append(ind)

        r_pr = OxmlElement("w:rPr")
        r_fonts = OxmlElement("w:rFonts")
        r_fonts.set(qn("w:ascii"), "Arial")
        r_fonts.set(qn("w:hAnsi"), "Arial")
        r_pr.append(r_fonts)

        for el in (start, num_fmt, lvl_text, lvl_jc, p_pr, r_pr):
            lvl.append(el)
        abstract_num.append(lvl)

    numbering_el.append(abstract_num)

    num = OxmlElement("w:num")
    num.set(qn("w:numId"), str(num_id))
    abstract_num_id = OxmlElement("w:abstractNumId")
    abstract_num_id.set(qn("w:val"), str(abstract_id))
    num.append(abstract_num_id)
    numbering_el.append(num)

    return num_id


def _set_bullet_level(paragraph, num_id: int, ilvl: int) -> None:
    """Attach paragraph to the given numId at the given ilvl (= outline level)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    p_pr = paragraph._p.get_or_add_pPr()
    num_pr = OxmlElement("w:numPr")
    ilvl_el = OxmlElement("w:ilvl")
    ilvl_el.set(qn("w:val"), str(ilvl))
    num_id_el = OxmlElement("w:numId")
    num_id_el.set(qn("w:val"), str(num_id))
    num_pr.append(ilvl_el)
    num_pr.append(num_id_el)
    p_pr.append(num_pr)


def outline_to_docx_bytes(outline: list[dict]) -> bytes:
    """
    Build a .docx that mirrors real GenMeet minutes formatting:
      level 0   → "Heading 2" style (new section)
      level 1   → bulleted paragraph, bold (topic)
      level 2+  → bulleted paragraph, normal weight, nested via ilvl

    This is the same representation Google Docs itself produces (heading
    style for section breaks, a single multilevel bullet list — via ilvl —
    for everything else), so a .docx generated here round-trips cleanly
    through read_outline_file() / _docx_paragraph_level() in Tab 2, and
    also reads correctly if the person edits it in Word/Google Docs first.
    """
    from docx import Document
    from docx.shared import Pt

    doc = Document()

    # ── Strip the default empty first paragraph ───────────────────────────
    for p in doc.paragraphs:
        p._element.getparent().remove(p._element)

    # ── Remove space-after on Normal style so lines sit flush ─────────────
    style = doc.styles["Normal"]
    style.paragraph_format.space_before = Pt(0)
    style.paragraph_format.space_after  = Pt(0)

    max_level = max((item["level"] for item in outline), default=0)
    num_id = _add_bullet_numbering(doc, max_level=max(max_level, 1))

    for item in outline:
        lvl  = item["level"]
        text = item["text"]

        if lvl == 0:
            doc.add_paragraph(text, style="Heading 2")
            continue

        para = doc.add_paragraph()
        para.paragraph_format.space_after = Pt(0)
        _set_bullet_level(para, num_id, lvl)

        run = para.add_run(text)
        if lvl == 1:
            run.bold = True

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _docx_paragraph_level(paragraph, last_level: int) -> int:
    """
    Determine a docx paragraph's outline level.

    Both Google Docs exports and our own outline_to_docx_bytes() output
    encode nesting the same way: a paragraph's list level (w:numPr/w:ilvl)
    *is* the outline level (0 = section, 1 = topic, 2+ = body). This is
    checked first and is authoritative — paragraph *style* is not a
    reliable signal on its own, since real minutes docs sometimes apply a
    Heading style to a sub-item purely for visual emphasis while still
    nesting it correctly via ilvl (e.g. a "Committee Updates" topic styled
    as Heading 2 but at ilvl=1, not a real new section).

    Falls back to:
      - level 0 if the paragraph uses a Heading style and has no list level
        at all (covers manually-typed headings with no bullets under them)
      - the previous paragraph's level, to gracefully handle stray
        unstyled/unbulleted paragraphs (e.g. blank formatting artifacts)
    """
    p_pr = paragraph._p.pPr
    if p_pr is not None and p_pr.numPr is not None and p_pr.numPr.ilvl is not None:
        return p_pr.numPr.ilvl.val

    style_name = (paragraph.style.name or "") if paragraph.style else ""
    if style_name.lower().startswith("heading") or style_name.lower() == "title":
        return 0

    return last_level


def read_outline_file(file) -> list[str]:
    """
    Read lines from an uploaded outline file, re-encoded as tab-indented
    text (one leading \t per outline level) so downstream parsing always
    goes through the unambiguous tab-count path in _parse_level(), never
    the Roman/alpha/numeric prefix-guessing fallback.

    .docx: level comes from each paragraph's real Word/Google Docs list
    level and heading style (see _docx_paragraph_level) — this correctly
    handles actual minutes docs (headings + native bulleted lists), not
    just outline.docx files produced by Tab 1 itself.

    .txt: passed through unchanged; _parse_level() handles tabs/prefixes.
    """
    name = getattr(file, "name", "")
    if name.lower().endswith(".docx"):
        from docx import Document
        doc = Document(file)
        lines: list[str] = []
        last_level = 0
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            level = _docx_paragraph_level(p, last_level)
            last_level = level
            lines.append("\t" * level + text)
        return lines
    return file.read().decode("utf-8").splitlines()


# ══════════════════════════════════════════════════════════════════════════════
# GOOGLE DOCS HELPERS
# ══════════════════════════════════════════════════════════════════════════════

DOCS_SCOPE  = "https://www.googleapis.com/auth/documents"
DRIVE_SCOPE = "https://www.googleapis.com/auth/drive"
INDENT_PT   = 18  # points per level

def _get_service_account_creds(scopes: list[str]):
    """Build credentials from Streamlit secrets (service account JSON)."""
    try:
        from google.oauth2 import service_account as sa
        info = dict(st.secrets["gcp_service_account"])
        # private_key newlines may be escaped in TOML
        if "\\n" in info.get("private_key", ""):
            info["private_key"] = info["private_key"].replace("\\n", "\n")
        return sa.Credentials.from_service_account_info(info, scopes=scopes)
    except Exception as e:
        st.error(f"❌ Could not load Google credentials from secrets: {e}")
        return None


def gdoc_write_outline(outline: list[dict], title: str, share_email: Optional[str] = None) -> Optional[str]:
    """Create a Google Doc with the formatted outline. Returns doc URL or None."""
    from googleapiclient.discovery import build as gbuild

    scopes = [DOCS_SCOPE, DRIVE_SCOPE] if share_email else [DOCS_SCOPE]
    creds = _get_service_account_creds(scopes)
    if not creds:
        return None

    docs_svc  = gbuild("docs",  "v1", credentials=creds)
    drive_svc = gbuild("drive", "v3", credentials=creds) if share_email else None

    # Create doc
    doc    = docs_svc.documents().create(body={"title": title}).execute()
    doc_id = doc["documentId"]

    # Insert all text — plain text only, no tabs; API handles indentation
    full_text = "\n".join(item["text"] for item in outline)
    docs_svc.documents().batchUpdate(
        documentId=doc_id,
        body={"requests": [{"insertText": {"location": {"index": 1}, "text": full_text}}]},
    ).execute()

    # Apply formatting per line
    fmt_requests = []
    char_idx = 1
    for i, item in enumerate(outline):
        ln  = len(item["text"])
        lvl = item["level"]
        start, end = char_idx, char_idx + ln

        fmt_requests.append({
            "updateParagraphStyle": {
                "range": {"startIndex": start, "endIndex": end},
                "paragraphStyle": {
                    "indentStart": {"magnitude": lvl * INDENT_PT, "unit": "PT"},
                    "spaceAbove":  {"magnitude": 4 if lvl == 0 else 0, "unit": "PT"},
                },
                "fields": "indentStart,spaceAbove",
            }
        })
        if lvl == 0:
            fmt_requests.append({
                "updateTextStyle": {
                    "range": {"startIndex": start, "endIndex": end},
                    "textStyle": {"bold": True, "fontSize": {"magnitude": 13, "unit": "PT"}},
                    "fields": "bold,fontSize",
                }
            })
        elif lvl == 1:
            fmt_requests.append({
                "updateTextStyle": {
                    "range": {"startIndex": start, "endIndex": end},
                    "textStyle": {"bold": True},
                    "fields": "bold",
                }
            })
        char_idx = end + (1 if i < len(outline) - 1 else 0)

    if fmt_requests:
        docs_svc.documents().batchUpdate(
            documentId=doc_id, body={"requests": fmt_requests}
        ).execute()

    # Share with user's email if provided
    if drive_svc and share_email:
        drive_svc.permissions().create(
            fileId=doc_id,
            body={"type": "user", "role": "writer", "emailAddress": share_email},
            sendNotificationEmail=False,
        ).execute()

    return f"https://docs.google.com/document/d/{doc_id}/edit"


def gdoc_read_lines(doc_id: str) -> list[str]:
    """Read non-empty lines from a Google Doc."""
    from googleapiclient.discovery import build as gbuild
    creds = _get_service_account_creds([DOCS_SCOPE])
    if not creds:
        return []
    svc = gbuild("docs", "v1", credentials=creds)
    doc = svc.documents().get(documentId=doc_id).execute()
    lines = []
    for element in doc.get("body", {}).get("content", []):
        para = element.get("paragraph")
        if not para:
            continue
        text = "".join(
            el["textRun"]["content"]
            for el in para.get("elements", [])
            if "textRun" in el
        ).strip()
        if text:
            lines.append(text)
    return lines


# ══════════════════════════════════════════════════════════════════════════════
# CSV GENERATION — GenMeet section-aware
# ══════════════════════════════════════════════════════════════════════════════

_ROMAN_RE = re.compile(r"^(M{0,4}(?:CM|CD|D?C{0,3})(?:XC|XL|L?X{0,3})(?:IX|IV|V?I{0,3}))\.\s+", re.I)
_ALPHA_RE = re.compile(r"^[A-Z]\.\s+")
_NUM_RE   = re.compile(r"^\d+\.\s+")
_LOWER_RE = re.compile(r"^[a-z]\)\s+")
_PAREN_RE = re.compile(r"^\(\d+\)\s+")

# Fixed section prefixes with their keyword triggers (case-insensitive substring match).
# Order here defines the column order in the output CSV.
SECTION_KEYWORDS: dict[str, list[str]] = {
    "Mem":  ["mem"],
    "Exec": ["exec"],
    "Info": ["info"],
    "Acti": ["acti"],
    "Exte": ["exte"],
    "Fin":  ["fin"],
    "Acad": ["acad"],
    "Rex":  ["rex"],
}

def _parse_level(line: str) -> tuple[int, str]:
    """
    Detect outline level and return (level, bare_text).

    Priority order:
      1. Leading tabs   → level = tab count  (Tab 1 docx/session output)
      2. Known prefix   → level from regex   (manually maintained Google Doc)
      3. Non-blank, no tabs, no prefix → level 0 (top-level heading in tab format)
      4. Blank          → -1, skip

    Rule 3 is safe because in the tab-based format level-0 items genuinely
    have no tabs; in the prefix-based Google Doc format any stray unprefixed
    line gets tested against section keywords and silently skipped if it
    doesn't match.
    """
    # ── Tab-indented (Tab 1 output) ───────────────────────────────────────────
    stripped = line.lstrip("\t")
    if len(stripped) < len(line):                 # had at least one leading tab
        return len(line) - len(stripped), stripped.strip()

    # ── Prefix-based (Google Doc source) ──────────────────────────────────────
    s = line.lstrip()
    if not s:
        return -1, s                              # blank line → skip
    if _ROMAN_RE.match(s): return 0, _ROMAN_RE.sub("", s).strip()
    if _ALPHA_RE.match(s): return 1, _ALPHA_RE.sub("", s).strip()
    if _NUM_RE.match(s):   return 2, _NUM_RE.sub("",   s).strip()
    if _LOWER_RE.match(s): return 3, _LOWER_RE.sub("", s).strip()
    if _PAREN_RE.match(s): return 4, _PAREN_RE.sub("", s).strip()

    # ── No tabs, no prefix, non-blank → top-level heading ────────────────────
    return 0, s

def _match_section_prefix(title: str) -> str | None:
    """Return the section prefix whose keyword appears in the title, or None."""
    t = title.lower()
    for prefix, keywords in SECTION_KEYWORDS.items():
        if any(kw in t for kw in keywords):
            return prefix
    return None

def _format_body(lines: list[str]) -> str:
    """Join body lines as a bulleted list. Always prefixed with • for consistency."""
    return "\n".join(f"• {ln}" for ln in lines)

def parse_meeting_sections(lines: list[str]) -> dict[str, dict]:
    """
    Parse outline lines into a dict keyed by section prefix:
        {
          "Mem":  {"heading": "...", "topics": [("Topic text", "• body\n• body"), ...]},
          "Exec": {"heading": "...", "topics": [...]},
          ...
        }

    Levels A. (1) → Topic
    Levels 1. / a) / (1) nested under a topic → all collected as body bullet lines.
    """
    sections: dict[str, dict] = {}
    current_prefix: str | None = None
    current_topic:  str | None = None
    current_bodies: list[str]  = []

    def _flush_topic():
        """Save the buffered topic+bodies into the current section."""
        if current_prefix and current_topic is not None:
            body_text = _format_body(current_bodies)
            sections[current_prefix]["topics"].append((current_topic, body_text))

    for line in lines:
        lvl, text = _parse_level(line)
        if lvl < 0:
            continue

        if lvl == 0:                                   # Roman numeral → new section
            _flush_topic()
            current_topic  = None
            current_bodies = []
            prefix = _match_section_prefix(text)
            current_prefix = prefix
            if prefix:
                sections[prefix] = {"heading": text, "topics": []}

        elif lvl == 1 and current_prefix:              # A. → new topic within section
            _flush_topic()
            current_topic  = text
            current_bodies = []

        elif lvl >= 2 and current_prefix and current_topic is not None:
            # 1. / a) / (1) → body content for the current topic
            current_bodies.append(text)

    _flush_topic()                                     # flush final topic
    return sections


def generate_meeting_csv(sections: dict[str, dict]) -> tuple[bytes, list[str]]:
    """
    Produce a single-row CSV where columns are:
        {prefix} Heading | {prefix} Topic 1 | {prefix} Body 1 | {prefix} Topic 2 | …
    Column count per section is driven by how many topics that section actually has.
    Sections are ordered by SECTION_KEYWORDS definition.
    """
    headers: list[str] = []
    for prefix in SECTION_KEYWORDS:
        data = sections.get(prefix, {"heading": "", "topics": []})
        headers.append(f"{prefix} Heading")
        for i in range(1, len(data["topics"]) + 1):
            headers.append(f"{prefix} Topic {i}")
            headers.append(f"{prefix} Body {i}")

    row: dict[str, str] = {}
    for prefix in SECTION_KEYWORDS:
        data = sections.get(prefix, {"heading": "", "topics": []})
        row[f"{prefix} Heading"] = data["heading"]
        for i, (topic, body) in enumerate(data["topics"], 1):
            row[f"{prefix} Topic {i}"] = topic
            row[f"{prefix} Body {i}"]  = body

    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=headers, extrasaction="ignore")
    writer.writeheader()
    writer.writerow(row)
    return buf.getvalue().encode("utf-8"), headers


# ══════════════════════════════════════════════════════════════════════════════
# HEADER
# ══════════════════════════════════════════════════════════════════════════════

st.markdown("""
<div class="top-header">
  <h1>🎨 UP CURSOR GenMeet Highlights Automater</h1>
  <p>Extract slide text → format into an outline → push to Google Docs → generate Canva Bulk Create CSV</p>
  <br>
  <span class="badge">100% Free</span>
  <span class="badge">No Canva API needed</span>
  <span class="badge">Google Docs</span>
  <span class="badge">Bulk Create</span>
</div>
""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════════════════════
# TABS
# ══════════════════════════════════════════════════════════════════════════════

tab1, tab2, tab3 = st.tabs([
    "📄  Step 1 — PPTX → Google Doc",
    "📊  Step 2 — Google Doc → Canva CSV",
    "🗂️  Attendance",
])


# ─────────────────────────────────────────────────────────────────────────────
#  TAB 1
# ─────────────────────────────────────────────────────────────────────────────
with tab1:
    st.markdown('<div class="step-pill">① Upload your Canva export and build the outline</div>', unsafe_allow_html=True)

    col_up, col_opts = st.columns([1.4, 1], gap="large")

    with col_up:
        st.markdown("#### Upload Canva PPTX")
        pptx_file = st.file_uploader(
            "Export from Canva → Share → Download → PowerPoint (.pptx)",
            type=["pptx"],
            key="pptx_upload",
        )

    with col_opts:
        st.markdown("#### Options")
        doc_title = st.text_input("Google Doc title", value="Presentation Outline", key="doc_title")
        google_email = st.text_input(
            "Share doc with (your Gmail)",
            placeholder="you@gmail.com",
            help="The created doc will be shared with this address so it appears in your Drive.",
            key="google_email",
        )

        gdocs_enabled = "gcp_service_account" in st.secrets
        if gdocs_enabled:
            st.markdown('<div class="success-box">✅ Google service account detected — Google Docs push is enabled.</div>', unsafe_allow_html=True)
        else:
            st.markdown('<div class="warn-box">⚠️ No Google credentials found in secrets. You can still download the outline as a text file.</div>', unsafe_allow_html=True)

    # ── Process PPTX ─────────────────────────────────────────────────────────
    if pptx_file:
        with st.spinner("Extracting slides …"):
            pptx_bytes = pptx_file.read()
            slides  = extract_slides(pptx_bytes)
            outline = build_numbered_outline(slides)
            # Save to session state so Tab 2 can pick it up
            st.session_state["outline"] = outline
            st.session_state["outline_lines"] = outline_to_plain_text(outline).split("\n")

        st.markdown(f"**{len(slides)} slides → {len(outline)} outline items**")
        st.markdown("#### Outline Preview")
        st.markdown(outline_to_html_preview(outline), unsafe_allow_html=True)

        st.markdown("---")
        dl_col, gdoc_col = st.columns(2, gap="medium")

        # ── Download as text ─────────────────────────────────────────────────
        with dl_col:
            st.markdown("##### 📥 Option A — Download as Word doc")
            docx_bytes = outline_to_docx_bytes(outline)
            st.download_button(
                label="⬇️  Download outline.docx",
                data=docx_bytes,
                file_name="outline.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
            )
            st.caption("Tab-indented outline — open in Word or Google Docs, then copy-paste. Upload in Tab 2 to generate the Canva CSV.")

        # ── Push to Google Doc ────────────────────────────────────────────────
        with gdoc_col:
            st.markdown("##### ☁️  Option B — Push to Google Doc")
            if gdocs_enabled:
                if st.button("🚀  Create Google Doc", use_container_width=True, type="primary"):
                    if not google_email:
                        st.warning("Enter your Gmail address above to share the doc with yourself.")
                    else:
                        with st.spinner("Creating Google Doc …"):
                            url = gdoc_write_outline(outline, doc_title, google_email)
                        if url:
                            doc_id = url.split("/d/")[1].split("/")[0]
                            st.session_state["last_doc_id"] = doc_id
                            st.markdown(f'<div class="success-box">✅ Doc created! <a href="{url}" target="_blank">Open in Google Docs ↗</a><br><small>Doc ID: <code>{doc_id}</code> (saved for Tab 2)</small></div>', unsafe_allow_html=True)
            else:
                st.button("🚀  Create Google Doc", use_container_width=True, disabled=True)
                st.caption("Add a service account to Streamlit secrets to enable this. See README.")
    else:
        st.markdown("""
<div class="info-box">
  👆 Upload a <strong>.pptx</strong> file exported from Canva to get started.<br><br>
  In Canva: <strong>Share → Download → Microsoft PowerPoint (.pptx)</strong>
</div>
""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────────────────
#  TAB 2
# ─────────────────────────────────────────────────────────────────────────────
with tab2:
    st.markdown('<div class="step-pill">② Generate the Canva Bulk Create CSV from your outline</div>', unsafe_allow_html=True)

    # ── Source selection ──────────────────────────────────────────────────────
    has_session_outline = "outline_lines" in st.session_state
    gdocs_enabled2      = "gcp_service_account" in st.secrets

    source_options = ["From Step 1 (this session)"] if has_session_outline else []
    source_options += ["Upload outline file (.docx or .txt)", "Google Doc ID"]
    source = st.radio("Outline source", source_options, horizontal=True)

    lines: list[str] = []

    if source == "From Step 1 (this session)" and has_session_outline:
        lines = st.session_state["outline_lines"]
        st.markdown(f'<div class="success-box">✅ Using outline from Step 1 — {len(lines)} lines loaded.</div>', unsafe_allow_html=True)

    elif source == "Upload outline file (.docx or .txt)":
        outline_file = st.file_uploader(
            "Upload the outline.docx (or .txt) downloaded in Step 1",
            type=["docx", "txt"],
            key="outline_upload",
        )
        if outline_file:
            lines = read_outline_file(outline_file)
            st.success(f"Loaded {len(lines)} lines from {outline_file.name}.")

    elif source == "Google Doc ID":
        default_id = st.session_state.get("last_doc_id", "")
        doc_id_input = st.text_input(
            "Google Doc ID",
            value=default_id,
            placeholder="Paste the document ID from the URL",
            help="Found in the URL: docs.google.com/document/d/**THIS_PART**/edit",
        )
        if doc_id_input:
            if not gdocs_enabled2:
                st.markdown('<div class="warn-box">⚠️ Google credentials not configured — cannot read Google Docs directly. Download the outline.docx from Step 1 and upload it instead.</div>', unsafe_allow_html=True)
            else:
                if st.button("📥 Load from Google Doc"):
                    with st.spinner("Reading Google Doc …"):
                        lines = gdoc_read_lines(doc_id_input.strip())
                    if lines:
                        st.session_state["loaded_lines"] = lines
                        st.success(f"Loaded {len(lines)} lines from Google Doc.")
                    else:
                        st.error("Could not read the doc. Make sure the service account has access.")
                if "loaded_lines" in st.session_state:
                    lines = st.session_state["loaded_lines"]

    # ── Generate CSV ──────────────────────────────────────────────────────────
    if lines:
        sections = parse_meeting_sections(lines)

        unmatched = [s for s in SECTION_KEYWORDS if s not in sections]
        if unmatched:
            st.warning(f"Could not match sections: {', '.join(unmatched)}. Check that the outline headings contain the expected keywords.")

        if not sections:
            st.warning("No sections detected. Make sure the outline uses Roman numeral headings that contain section keywords (Mem, Exec, Info, Acti, Exte, Fin, Acad, Rex).")
        else:
            total_topics = sum(len(d["topics"]) for d in sections.values())
            st.markdown(f"**{len(sections)} sections · {total_topics} total topics parsed**")

            csv_bytes, headers = generate_meeting_csv(sections)

            # ── Section preview cards ─────────────────────────────────────────
            st.markdown("#### Section Preview")
            for prefix in SECTION_KEYWORDS:
                if prefix not in sections:
                    continue
                data = sections[prefix]
                with st.expander(f"**{prefix}** — {data['heading']}  ({len(data['topics'])} topics)", expanded=False):
                    for i, (topic, body) in enumerate(data["topics"], 1):
                        st.markdown(f"**Topic {i}:** {topic}")
                        st.code(body, language=None)

            # ── CSV preview ───────────────────────────────────────────────────
            st.markdown("#### CSV Preview")
            import pandas as pd
            df = pd.read_csv(io.BytesIO(csv_bytes))
            st.dataframe(df.T.rename(columns={0: "Value"}), use_container_width=True, height=320)
            st.caption("Transposed for readability — one column per section field, value on the right.")

            st.download_button(
                label="⬇️  Download canva_bulk_create.csv",
                data=csv_bytes,
                file_name="canva_bulk_create.csv",
                mime="text/csv",
                type="primary",
            )

            # ── Canva field name table ────────────────────────────────────────
            st.markdown("---")
            st.markdown("#### 🎨 Canva field names to connect in your template")
            st.markdown("""
<div class="info-box">
  For each textbox in your Canva template, click it → <strong>Connect data</strong> → type the field name exactly as shown below.
  Then go to <strong>Apps → Bulk Create → Upload CSV</strong>.
</div>
""", unsafe_allow_html=True)

            rows_html = ""
            for prefix in SECTION_KEYWORDS:
                if prefix not in sections:
                    continue
                data    = sections[prefix]
                n_topics = len(data["topics"])
                rows_html += f"<tr><td colspan='2' style='background:#f5f3ff;font-weight:600;color:#4338ca;padding:.5rem .8rem'>{prefix} — {data['heading']}</td></tr>"
                rows_html += f"<tr><td>Section heading textbox</td><td><span class='code-chip'>{prefix} Heading</span></td></tr>"
                for i in range(1, n_topics + 1):
                    rows_html += f"<tr><td>Topic {i} textbox</td><td><span class='code-chip'>{prefix} Topic {i}</span></td></tr>"
                    rows_html += f"<tr><td>Body {i} textbox</td><td><span class='code-chip'>{prefix} Body {i}</span></td></tr>"

            st.markdown(f"""
<table class="field-table">
  <thead><tr><th>Textbox purpose</th><th>Field name to enter in Canva</th></tr></thead>
  <tbody>{rows_html}</tbody>
</table>
""", unsafe_allow_html=True)

    else:
        if not has_session_outline:
            st.markdown("""
<div class="info-box">
  Complete <strong>Step 1</strong> first, or upload an <code>outline.txt</code> file here directly.
</div>
""", unsafe_allow_html=True)


# ─────────────────────────────────────────────────────────────────────────────
#  TAB 3 — ATTENDANCE
# ─────────────────────────────────────────────────────────────────────────────
with tab3:
    st.markdown('<div class="step-pill">③ Upload F2F and/or Zoom attendance files to consolidate</div>', unsafe_allow_html=True)

    col_f2f, col_zoom = st.columns(2, gap="large")

    with col_f2f:
        st.markdown("#### F2F Attendance")
        st.caption("CSV or Excel — columns: Surname, Nickname, Entry Time, Status")
        f2f_file = st.file_uploader("Upload F2F file", type=["csv", "xlsx", "xls"], key="f2f_upload")

    with col_zoom:
        st.markdown("#### Zoom Attendance")
        st.caption("CSV — columns: Nickname, Duration, Status")
        zoom_file = st.file_uploader("Upload Zoom file", type=["csv"], key="zoom_upload")

    if not f2f_file and not zoom_file:
        st.markdown("""
<div class="info-box">
  Upload at least one file above.<br>
  Both sources are optional — uploading only one will use that source's status directly.
</div>
""", unsafe_allow_html=True)
    else:
        try:
            # ── Load & validate ───────────────────────────────────────────────
            f2f_df  = att.load_f2f(f2f_file)   if f2f_file  else None
            zoom_df = att.load_zoom(zoom_file)  if zoom_file else None

            # ── Source previews ───────────────────────────────────────────────
            with st.expander("Source data preview", expanded=False):
                if f2f_df is not None:
                    st.markdown("**F2F**")
                    st.dataframe(
                        f2f_df.drop(columns=["_nick_key"]),
                        use_container_width=True, height=200,
                    )
                if zoom_df is not None:
                    st.markdown("**Zoom** (deduplicated — reconnections merged)")
                    st.dataframe(
                        zoom_df.drop(columns=["_nick_key"]),
                        use_container_width=True, height=200,
                    )

            # ── Consolidate ───────────────────────────────────────────────────
            result = att.consolidate(f2f_df, zoom_df)

            counts = result["Final Status"].value_counts()
            n_att  = counts.get("Attendee", 0)
            n_late = counts.get("Late",     0)
            n_abs  = counts.get("Absent",   0)

            # Summary metric chips
            st.markdown(f"""
<div style="display:flex;gap:1rem;margin:1rem 0">
  <div class="card" style="flex:1;text-align:center">
    <div style="font-size:1.8rem;font-weight:700;color:#15803d">{n_att}</div>
    <div style="font-size:.8rem;color:#6b7280">Attendees</div>
  </div>
  <div class="card" style="flex:1;text-align:center">
    <div style="font-size:1.8rem;font-weight:700;color:#854d0e">{n_late}</div>
    <div style="font-size:.8rem;color:#6b7280">Late</div>
  </div>
  <div class="card" style="flex:1;text-align:center">
    <div style="font-size:1.8rem;font-weight:700;color:#991b1b">{n_abs}</div>
    <div style="font-size:.8rem;color:#6b7280">Absent</div>
  </div>
  <div class="card" style="flex:1;text-align:center">
    <div style="font-size:1.8rem;font-weight:700;color:#1e40af">{len(result)}</div>
    <div style="font-size:.8rem;color:#6b7280">Total</div>
  </div>
</div>
""", unsafe_allow_html=True)

            # ── Colour-coded result table ─────────────────────────────────────
            STATUS_COLORS = {
                "Attendee": "background-color:#dcfce7;color:#15803d;font-weight:600",
                "Late":     "background-color:#fef9c3;color:#854d0e;font-weight:600",
                "Absent":   "background-color:#fee2e2;color:#991b1b;font-weight:600",
            }

            rows_html = ""
            for _, row in result.iterrows():
                style = STATUS_COLORS.get(row["Final Status"], "")
                rows_html += (
                    f"<tr>"
                    f"<td style='padding:.45rem .8rem;border-bottom:1px solid #f0f0f0'>{row['Nickname']}</td>"
                    f"<td style='padding:.45rem .8rem;border-bottom:1px solid #f0f0f0;{style}'>{row['Final Status']}</td>"
                    f"</tr>"
                )

            st.markdown(f"""
<table class="field-table" style="margin-top:.5rem">
  <thead><tr><th>Nickname</th><th>Final Status</th></tr></thead>
  <tbody>{rows_html}</tbody>
</table>
""", unsafe_allow_html=True)

            # ── Download ──────────────────────────────────────────────────────
            csv_out = result.to_csv(index=False).encode("utf-8")
            st.download_button(
                label="⬇️  Download consolidated_attendance.csv",
                data=csv_out,
                file_name="consolidated_attendance.csv",
                mime="text/csv",
                type="primary",
            )

        except ValueError as e:
            st.error(f"Column error: {e}")
        except Exception as e:
            st.error(f"Unexpected error: {e}")
